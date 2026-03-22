"""
Builds a presentation deck for the Expense Tracker ML project using existing figures.
Requires: python-pptx (install with `pip install python-pptx` when online).
Output: ExpenseTrackerML.pptx in the repo root.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


FIG = Path("figures")


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Expense Tracker ML"
    slide.placeholders[1].text = "Automated expense categorization + budget alerts\nTeam: <your names>\nDate: <today>"


def add_text_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    for i, line in enumerate(bullets):
        if i == 0:
            body.text = line
        else:
            body.add_paragraph().text = line


def add_image_slide(prs, title, image_path, height_inches=5.0):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.3)
    height = Inches(height_inches)
    slide.shapes.add_picture(str(image_path), left, top, height=height)


def build_deck():
    prs = Presentation()

    add_title_slide(prs)

    add_text_slide(
        prs,
        "Problem & Goal",
        [
            "Manual expense categorization is slow and inconsistent.",
            "Goal: fast, accurate category prediction + clear spending insights.",
            "Data: spending_patterns_detailed.csv (10k rows, 9 columns).",
        ],
    )

    add_image_slide(prs, "Class Distribution", FIG / "class_distribution.png")
    add_image_slide(prs, "Architecture", FIG / "architecture_diagram.png", height_inches=4.5)

    add_text_slide(
        prs,
        "Feature Pipeline",
        [
            "TF-IDF on Item text (max_features=3000).",
            "StandardScaler on Total Spent.",
            "One-Hot on Payment Method & Location.",
            "Rare categories (<3) folded into 'Other'; stratified 80/20 split.",
        ],
    )

    add_image_slide(prs, "Model Comparison (Metrics Table)", FIG / "model_metric_table.png", height_inches=4.5)
    add_image_slide(prs, "ROC Comparison — All Models", FIG / "roc_curve_all_models.png", height_inches=4.5)

    add_image_slide(prs, "Precision–Recall (LogReg)", FIG / "precision_recall_logreg.png", height_inches=4.5)
    add_image_slide(prs, "Confusion Matrix (LogReg)", FIG / "confusion_matrix_logreg.png", height_inches=4.5)
    add_image_slide(prs, "Top Misclassifications", FIG / "top_confusions_logreg.png", height_inches=4.0)
    add_image_slide(prs, "Feature Importance (LogReg)", FIG / "feature_importance_logreg.png", height_inches=4.5)
    add_image_slide(prs, "Learning Curve (LogReg)", FIG / "learning_curve_logreg.png", height_inches=4.0)

    add_text_slide(
        prs,
        "Results & Next Steps",
        [
            "Best overall: Logistic Regression (balanced class weights).",
            "Confidence available via predict_proba for UI display.",
            "Next: hyperparameter sweep, incremental retraining, mobile UI.",
        ],
    )

    out = Path("ExpenseTrackerML.pptx")
    prs.save(out)
    print(f"Saved deck to {out.resolve()}")


if __name__ == "__main__":
    build_deck()
